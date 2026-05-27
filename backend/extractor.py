# -*- coding: utf-8 -*-
# 텍스트 추출기 - 문서 유형별 최적 추출
import re
import zlib
import struct
import zipfile
from pathlib import Path
from typing import Optional

import os
os.environ['JAVA_TOOL_OPTIONS'] = '-Djava.awt.headless=true'

try:
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = \
        r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except ImportError:
    pass

SUPPORTED = {".pdf", ".hwp", ".hwpx", ".pptx", ".docx", ".xlsx",
             ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"}

IMAGE_EXT = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"}
OCR_MIN_CHARS = 50  # 이 이하면 품질 낮다고 판단 → LLM Vision fallback

# 파일명 strong 키워드 (이것만 있어도 분류 가능)
_FNAME_STRONG = [
    "계약서","협약서","견적서","세금계산서","정산서","발주서","전자세금계산서",
    "수행계획서","사업계획서","기획서","기획안","실행계획서","추진계획서",
    "발표자료","프레젠테이션","IR자료","소개자료",
    "공고문","모집공고","입찰공고","채용공고","제안요청서","과업지시서",
    "사업자등록증","법인등기","납세증명","중소기업확인서","벤처기업확인서","등록증",
    "동의서","개인정보수집이용동의서","참여동의서","수행계획서",
    "최종보고서","중간보고서","결과보고서","완료보고서",
]

# 문서 유형별 추출 전략
#   chars: 최대 추출 글자 수
#   pages: PDF 최대 페이지 수
#   slides: PPTX 최대 슬라이드 수
_STRATEGY = {
    # 파일명으로 판단 가능 → 텍스트 최소화
    "fname_only":   {"chars": 0,    "pages": 0, "slides": 0},
    # 1페이지면 충분 (계약서, 인증서, 동의서)
    "first_page":   {"chars": 800,  "pages": 1, "slides": 1},
    # 제목+목차+1페이지 (계획서, 공고문)
    "title_toc":    {"chars": 1200, "pages": 2, "slides": 3},
    # 조사자료, 법령별표 → 더 읽어야 함
    "extended":     {"chars": 2000, "pages": 3, "slides": 5},
    # 판단 불가 → 기본값
    "default":      {"chars": 1500, "pages": 2, "slides": 3},
}


def _pick_strategy(filename: str) -> str:
    """
    파일명만으로 추출 전략 결정.
    비용/시간 최소화 핵심 함수.
    """
    name = filename.lower()

    # 파일명에 strong 키워드 → 텍스트 추출 불필요할 수 있음
    # (classifier에서 다시 확인하므로 최소만 추출)
    for kw in _FNAME_STRONG:
        if kw in filename:
            # 계약서/인증서 → 1페이지면 충분
            if any(k in filename for k in ["계약서","협약서","세금계산서","정산서","견적서",
                                             "사업자등록증","납세증명","확인서","등록증","동의서"]):
                return "first_page"
            # 계획서/공고 → 목차까지
            if any(k in filename for k in ["계획서","공고문","모집공고","제안요청서","과업지시"]):
                return "title_toc"
            return "title_toc"

    # 별표/고시/시행규칙 → 표가 핵심이라 더 읽어야 함
    if any(k in filename for k in ["별표","고시","시행규칙","기술기준"]):
        return "extended"

    # 발표자료(pptx) → 슬라이드 제목만
    if any(k in filename for k in ["발표","ppt","pt"]):
        return "first_page"

    return "default"


# ── 공통 정제 ─────────────────────────────────────────────────────────────────

def _clean(text: str, max_chars: int) -> Optional[str]:
    if not text:
        return None
    text = re.sub(r"[\u4E00-\u9FFF]{1,4}(?=\s|$)", "", text)  # HWP 노이즈
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", text)  # 제어문자
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text[:max_chars] if text else None


# ── PDF ──────────────────────────────────────────────────────────────────────

def _pdf(fp: Path, strategy: str) -> Optional[str]:
    import fitz
    s    = _STRATEGY[strategy]
    doc  = fitz.open(str(fp))
    n    = len(doc)

    if s["pages"] == 0 or n == 0:
        doc.close()
        return None

    first = doc[0].get_text().strip() if n > 0 else ""
    if len(first) < 50:
        # 스캔본 → OCR (시간 걸리지만 어쩔 수 없음)
        text = _ocr_pdf(doc, min(s["pages"], n))
    else:
        text = "\n".join(doc[i].get_text() for i in range(min(s["pages"], n)))

    doc.close()
    return _clean(text, s["chars"])


def _ocr_pdf(doc, n_pages: int) -> str:
    try:
        import pytesseract
        from PIL import Image
        import io
        texts = []
        for i in range(n_pages):
            pix  = doc[i].get_pixmap(dpi=150)
            img  = Image.open(io.BytesIO(pix.tobytes("png")))
            texts.append(pytesseract.image_to_string(img, lang="kor+eng"))
        return "\n".join(texts)
    except ImportError:
        return ""


# ── HWPX ─────────────────────────────────────────────────────────────────────

def _hwpx(fp: Path, strategy: str) -> Optional[str]:
    s = _STRATEGY[strategy]
    if s["pages"] == 0:
        return None
    with zipfile.ZipFile(str(fp)) as z:
        names = z.namelist()
        if "Preview/PrvText.txt" in names:
            text = z.read("Preview/PrvText.txt").decode("utf-8", errors="ignore")
            if text.strip():
                return _clean(text, s["chars"])
        sections = sorted(n for n in names if re.match(r"Contents/section\d+\.xml", n))
        texts = []
        for sec in sections[:s["pages"]]:
            xml   = z.read(sec).decode("utf-8", errors="ignore")
            parts = re.findall(r"<(?:hp:)?t\b[^>]*>(.*?)</(?:hp:)?t>", xml, re.DOTALL)
            raw   = re.sub(r"<[^>]+>", " ", " ".join(parts))
            if raw.strip():
                texts.append(raw.strip())
        return _clean(" ".join(texts), s["chars"]) if texts else None


# ── HWP (OLE2 + zlib) ────────────────────────────────────────────────────────

def _hwp(fp: Path, strategy: str) -> Optional[str]:
    s = _STRATEGY[strategy]
    if s["pages"] == 0:
        return None
    try:
        import olefile
        ole = olefile.OleFileIO(str(fp))

        compressed = True
        if ole.exists("FileHeader"):
            hdr        = ole.openstream("FileHeader").read()
            flags      = struct.unpack_from("<I", hdr, 36)[0]
            compressed = bool(flags & 0x01)

        sections = sorted(
            s for s in ole.listdir()
            if len(s) == 2 and s[0] == "BodyText"
        )
        texts = []
        for path in sections[:_STRATEGY[strategy]["pages"]]:
            try:
                data = ole.openstream(path).read()
                body = zlib.decompress(data, -15) if compressed else data
                t    = _parse_hwp_records(body)
                if t:
                    texts.append(t)
            except Exception:
                pass

        ole.close()
        result = _clean(" ".join(texts), _STRATEGY[strategy]["chars"])
        return result if result else _hwp_fallback(fp, _STRATEGY[strategy]["chars"])
    except Exception:
        return _hwp_fallback(fp, _STRATEGY[strategy]["chars"])


def _parse_hwp_records(data: bytes) -> str:
    texts = []
    i     = 0
    while i < len(data) - 4:
        try:
            hdr    = struct.unpack_from("<I", data, i)[0]
            tag_id = hdr & 0x3FF
            size   = (hdr >> 20) & 0xFFF
            i     += 4
            if size == 0xFFF:
                if i + 4 > len(data): break
                size = struct.unpack_from("<I", data, i)[0]
                i   += 4
            if tag_id == 67 and 0 < size < 20000:
                chunk = data[i:i+size]
                text  = _decode_para(chunk)
                if text:
                    texts.append(text)
            i += size
        except Exception:
            i += 1
    return " ".join(texts)


def _decode_para(chunk: bytes) -> str:
    out = []
    j   = 0
    while j < len(chunk) - 1:
        code = struct.unpack_from("<H", chunk, j)[0]
        if code in (0x0D, 0x0A): out.append("\n")
        elif code == 0x09:        out.append(" ")
        elif 0x20 <= code <= 0xD7FF or 0xE000 <= code <= 0xFFFF:
            try: out.append(chr(code))
            except Exception: pass
        j += 2
    return "".join(out).strip()


def _hwp_fallback(fp: Path, max_chars: int) -> Optional[str]:
    with open(str(fp), "rb") as f:
        raw = f.read()
    pattern = re.compile(rb"(?:[\xAC-\xD7][\x00-\xFF]){3,}")
    results = []
    for m in pattern.finditer(raw):
        try:
            t = m.group().decode("utf-16-le", errors="ignore").strip()
            if len(t) > 2:
                results.append(t)
        except Exception:
            pass
    return _clean(" ".join(results), max_chars) if results else None


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _pptx(fp: Path, strategy: str) -> Optional[str]:
    s = _STRATEGY[strategy]
    if s["slides"] == 0:
        return None
    from pptx import Presentation
    prs   = Presentation(str(fp))
    texts = []
    for slide in list(prs.slides)[:s["slides"]]:
        # 제목 우선
        if slide.shapes.title and slide.shapes.title.text.strip():
            texts.append(slide.shapes.title.text.strip())
        # 나머지 텍스트
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.text.strip() not in texts:
                    texts.append(shape.text.strip())
    return _clean("\n".join(texts), s["chars"])


# ── DOCX / XLSX ──────────────────────────────────────────────────────────────

def _docx(fp: Path, strategy: str) -> Optional[str]:
    s = _STRATEGY[strategy]
    if s["pages"] == 0:
        return None
    from docx import Document
    doc  = Document(str(fp))
    n    = max(s["pages"] * 10, 20)  # 페이지당 약 10단락 가정
    return _clean("\n".join(p.text for p in doc.paragraphs[:n] if p.text.strip()), s["chars"])


def _xlsx(fp: Path, strategy: str) -> Optional[str]:
    s = _STRATEGY[strategy]
    if s["pages"] == 0:
        return None
    import openpyxl
    wb    = openpyxl.load_workbook(str(fp), read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(max_row=30, values_only=True):
            for cell in row:
                if cell:
                    texts.append(str(cell))
    return _clean(" ".join(texts), s["chars"])


# ── 통합 추출 함수 ────────────────────────────────────────────────────────────

# ── 이미지 (OCR) ──────────────────────────────────────────────────────────────

def _image(fp: Path, strategy: str) -> Optional[str]:
    """
    이미지 → OCR 텍스트 추출
    1. OpenCV 전처리 (흑백 + 노이즈제거 + 이진화)
    2. pytesseract OCR
    3. 추출 텍스트 50자 미만 → Claude Vision fallback
    """
    # Step 1: OpenCV 전처리 + pytesseract
    text = _ocr_with_preprocess(fp)

    # Step 2: 품질 낮으면 Claude Vision fallback
    if not text or len(text.strip()) < OCR_MIN_CHARS:
        print(f"  [OCR] 품질 낮음 ({len(text.strip()) if text else 0}자) → Vision fallback")
        text = _ocr_with_vision(fp)

    return _clean(text, _STRATEGY[strategy]["chars"]) if text else None


def _ocr_with_preprocess(fp: Path) -> Optional[str]:
    """OpenCV 전처리 후 pytesseract OCR."""
    try:
        import cv2
        import pytesseract
        import numpy as np

        # 이미지 로드
        img = cv2.imread(str(fp))
        if img is None:
            return None

        # 전처리: 흑백 변환
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # 노이즈 제거
        blur = cv2.GaussianBlur(gray, (3, 3), 0)

        # 이진화 (Otsu 방식 - 자동 임계값)
        _, thresh = cv2.threshold(
            blur, 0, 255,
            cv2.THRESH_BINARY + cv2.THRESH_OTSU
        )

        # 기울기 보정 시도
        thresh = _deskew(thresh)

        # OCR
        text = pytesseract.image_to_string(thresh, lang="kor+eng")
        return text

    except ImportError:
        # opencv 없으면 PIL로 fallback
        return _ocr_without_preprocess(fp)
    except Exception as e:
        print(f"  [OCR 오류] {e}")
        return None


def _deskew(img):
    """기울어진 이미지 보정."""
    try:
        import cv2
        import numpy as np
        coords = np.column_stack(np.where(img > 0))
        if len(coords) == 0:
            return img
        angle = cv2.minAreaRect(coords)[-1]
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle
        # 5도 이내만 보정 (너무 많이 돌리면 오히려 나빠짐)
        if abs(angle) > 5:
            return img
        h, w = img.shape[:2]
        center = (w // 2, h // 2)
        M = cv2.getRotationMatrix2D(center, angle, 1.0)
        rotated = cv2.warpAffine(
            img, M, (w, h),
            flags=cv2.INTER_CUBIC,
            borderMode=cv2.BORDER_REPLICATE
        )
        return rotated
    except Exception:
        return img


def _ocr_without_preprocess(fp: Path) -> Optional[str]:
    """opencv 없을 때 PIL + pytesseract."""
    try:
        import pytesseract
        from PIL import Image
        img  = Image.open(str(fp))
        text = pytesseract.image_to_string(img, lang="kor+eng")
        return text
    except ImportError:
        print("  [OCR] pytesseract/PIL 미설치")
        return None
    except Exception as e:
        print(f"  [OCR 오류] {e}")
        return None


def _ocr_with_vision(fp: Path) -> Optional[str]:
    """
    Claude Vision API fallback.
    pytesseract 결과가 50자 미만일 때만 호출.
    """
    try:
        import anthropic    
        import base64

        with open(str(fp), "rb") as f:
            img_data = base64.standard_b64encode(f.read()).decode("utf-8")

        # 확장자로 media_type 결정
        ext = fp.suffix.lower()
        media_map = {
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".png": "image/png",  ".gif":  "image/gif",
            ".bmp": "image/bmp",  ".tiff": "image/tiff",
        }
        media_type = media_map.get(ext, "image/jpeg")

        client = anthropic.Anthropic()
        resp   = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=500,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type":       "base64",
                            "media_type": media_type,
                            "data":       img_data,
                        },
                    },
                    {
                        "type": "text",
                        "text": "이 이미지에서 텍스트를 추출해주세요. 텍스트만 출력하고 다른 설명은 하지 마세요.",
                    },
                ],
            }],
        )
        return resp.content[0].text.strip()

    except Exception as e:
        print(f"  [Vision 오류] {e}")
        return None


def extract_text(filepath: Path) -> Optional[str]:
    """
    파일명으로 전략 결정 → 최소한만 추출.
    비용/시간 최소화.
    """
    ext      = filepath.suffix.lower()
    strategy = _pick_strategy(filepath.name)

    handlers = {
        ".pdf":  _pdf,
        ".hwpx": _hwpx,
        ".hwp":  _hwp,
        ".pptx": _pptx,
        ".docx": _docx,
        ".xlsx": _xlsx,
    }
    # 이미지 형식은 _image로 통합
    if ext in IMAGE_EXT:
        handlers[ext] = _image
    fn = handlers.get(ext)
    if not fn:
        return None
    try:
        return fn(filepath, strategy)
    except Exception as e:
        print(f"[추출실패] {filepath.name}: {e}")
        return None